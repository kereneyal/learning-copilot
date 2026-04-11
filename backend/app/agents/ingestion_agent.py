import os
from typing import Any, Dict, Optional

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation


class IngestionAgent:
    """last_pdf_meta is set after extract_text when file_type is pdf (for logging)."""

    def __init__(self):
        self.last_pdf_meta: Optional[Dict[str, Any]] = None

    def extract_text(self, file_path: str, file_type: str) -> str:
        file_type = file_type.lower()
        self.last_pdf_meta = None

        if file_type == "pdf":
            return self._extract_from_pdf(file_path)
        elif file_type == "docx":
            return self._extract_from_docx(file_path)
        elif file_type == "pptx":
            return self._extract_from_pptx(file_path)
        elif file_type == "txt":
            return self._extract_from_txt(file_path)
        else:
            return ""

    def detect_language(self, text: str) -> str:
        hebrew_chars = sum(1 for ch in text if '\u0590' <= ch <= '\u05FF')
        english_chars = sum(1 for ch in text if ('a' <= ch.lower() <= 'z'))

        if hebrew_chars > english_chars:
            return "he"
        return "en"

    def _extract_from_pdf(self, file_path: str) -> str:
        from app.services.pdf_ocr_service import extract_pdf_for_ingestion

        result = extract_pdf_for_ingestion(file_path)
        self.last_pdf_meta = {
            "provider": result.provider,
            "pages_processed": result.pages_processed,
            "ocr_used": result.ocr_used,
        }
        return result.text or ""

    def _extract_from_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        return "\n".join([p.text for p in doc.paragraphs]).strip()

    def _extract_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs).strip()

    def _extract_from_txt(self, file_path: str) -> str:
        for encoding in ("utf-8", "utf-8-sig", "cp1255", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding, errors="replace") as f:
                    return f.read().strip()
            except Exception:
                continue
        with open(file_path, "r", errors="replace") as f:
            return f.read().strip()
